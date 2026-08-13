#!/usr/bin/env python3
"""Import the shop-internal MIYU PPT photos into text-free consultation assets."""

from __future__ import annotations

import argparse
import io
from pathlib import Path

from PIL import Image, ImageOps, ImageStat
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PPT = Path('/Users/oeuvre/Desktop/미유/무드진단/★무드 세부 분류. 20260805.pptx')
OUTPUT = ROOT / 'assets' / 'diagnosis' / 'reference'
GROUPS = ('a', 'b', 'c', 'd')
# These names follow the PPT, not the app's D-1~3 order. The app resolves them
# by type name so Clear/Sharp/Charisma cannot be interchanged.
TYPE_VISUAL_SOURCES = {
    'fantasy': {'face': 15, 'comparison': 17, 'fashion': (19, 20)},
    'fruity': {'face': 22, 'comparison': 23, 'fashion': (25, 26)},
    'soda': {'face': 28, 'comparison': 29, 'fashion': (31, 32)},
    'romantic': {'face': 45, 'comparison': 46, 'fashion': (48, 49)},
    'soft': {'face': 51, 'comparison': 52, 'fashion': (54, 55)},
    'elegance': {'face': 56, 'comparison': 58, 'fashion': (60, 61)},
    'vintage': {'face': 73, 'comparison': 74, 'fashion': (76, 77)},
    'refined': {'face': 79, 'comparison': 80, 'fashion': (82, 83)},
    'deep-chic': {'face': 85, 'comparison': 86, 'fashion': (88, 89)},
    'clear': {'face': 100, 'comparison': 101, 'fashion': (103, 104)},
    'sharp': {'face': 106, 'comparison': 107, 'fashion': (109, 110)},
    'charisma': {'face': 112, 'comparison': 113, 'fashion': (115, 116)},
}
FEMALE_HAIR_RECOMMENDED_SLIDES = {
    'a': (7, 8, 9), 'b': (37, 38, 39, 40), 'c': (66, 67, 68), 'd': (94, 95, 96)
}
FEMALE_HAIR_AVOID_SLIDES = {'a': (11,)}
MALE_HAIR_RECOMMENDED_SLIDES = {'a': (13,), 'b': (42,), 'c': (70,), 'd': (98,)}
MALE_HAIR_AVOID_SLIDES = {'a': (14,), 'b': (43,), 'c': (71,), 'd': (99,)}


def open_picture(shape):
    image = Image.open(io.BytesIO(shape.image.blob))
    return ImageOps.exif_transpose(image).convert('RGB')


def jpg(image: Image.Image, destination: Path, size=(720, 900)):
    destination.parent.mkdir(parents=True, exist_ok=True)
    image.thumbnail(size, Image.Resampling.LANCZOS)
    image.save(destination, 'JPEG', quality=78, optimize=True, progressive=True)


def picture_shapes(slide):
    return [shape for shape in slide.shapes if shape.shape_type == 13]


def intro_portrait(image: Image.Image):
    """Slide 2 stores each group card as one image; keep only its face reference."""
    width, height = image.size
    return image.crop((int(width * 0.64), int(height * 0.14), int(width * 0.98), int(height * 0.91)))


def sampled(images, limit=12):
    if len(images) <= limit:
        return images
    return [images[round(index * (len(images) - 1) / (limit - 1))] for index in range(limit)]


def is_blank_placeholder(image: Image.Image):
    """Ignore tiny, plain-white PPT spacer shapes that are not reference photos."""
    stat = ImageStat.Stat(image.resize((16, 16)))
    return min(stat.mean) >= 250 and max(stat.var) <= 1


def balanced_row_counts(count: int):
    """Split a board into balanced rows so the last row has no empty cells."""
    if count < 1:
        return []
    columns = min(4, max(1, (count + 1) // 2))
    rows = (count + columns - 1) // columns
    base, extra = divmod(count, rows)
    return [base + 1 if index < extra else base for index in range(rows)]


def collage_images(images, destination: Path, limit: int = 12):
    pictures = [image for image in sampled(images, limit) if not is_blank_placeholder(image)]
    if not pictures:
        raise SystemExit(f'No pictures available for {destination}')
    row_counts = balanced_row_counts(len(pictures))
    canvas_width, cell_height, gap = 900, 238, 10
    canvas = Image.new(
        'RGB',
        (canvas_width, len(row_counts) * cell_height + (len(row_counts) + 1) * gap),
        '#f5f3ef'
    )
    image_index = 0
    for row_index, row_count in enumerate(row_counts):
        cell_width = (canvas_width - (row_count + 1) * gap) // row_count
        y = gap + row_index * (cell_height + gap)
        for column_index in range(row_count):
            image = pictures[image_index]
            copy = ImageOps.contain(image, (cell_width, cell_height), Image.Resampling.LANCZOS)
            x = gap + column_index * (cell_width + gap) + (cell_width - copy.width) // 2
            canvas.paste(copy, (x, y + (cell_height - copy.height) // 2))
            image_index += 1
    jpg(canvas, destination, (720, 900))


def card_images(images, destination: Path, count: int = 3):
    """Export distinct source photos for the horizontally scrollable explanation cards."""
    pictures = [image for image in images if not is_blank_placeholder(image)]
    if not pictures:
        raise SystemExit(f'No pictures available for {destination}')
    selected = sampled(pictures, count)
    # Source slides normally contain at least three photos. If a legacy slide
    # has fewer, preserve its real source image instead of inventing a new one.
    while len(selected) < count:
        selected.append(selected[-1])
    for index, image in enumerate(selected[:count], start=1):
        jpg(image, destination / f'{index}.jpg', (420, 560))


def slide_images(presentation, numbers):
    return [open_picture(shape) for number in numbers for shape in picture_shapes(presentation.slides[number - 1])]


def comparison_images(presentation, slide_number, side):
    shapes = picture_shapes(presentation.slides[slide_number - 1])
    midpoint = 7_500_000
    selected = [shape for shape in shapes if (shape.left < midpoint) == (side == 'recommended')]
    if not selected:
        raise SystemExit(f'Missing {side} comparison pictures on slide {slide_number}')
    return [open_picture(shape) for shape in selected]


def import_assets(ppt_path: Path):
    if not ppt_path.is_file():
        raise SystemExit(f'Missing MIYU reference PPT: {ppt_path}')
    presentation = Presentation(ppt_path)

    intro_pictures = picture_shapes(presentation.slides[1])
    if len(intro_pictures) != 4:
        raise SystemExit(f'Expected 4 intro faces on slide 2, found {len(intro_pictures)}')
    for group, shape in zip(GROUPS, intro_pictures):
        jpg(intro_portrait(open_picture(shape)), OUTPUT / 'intro' / f'{group}.jpg')

    for key, source in TYPE_VISUAL_SOURCES.items():
        pictures = picture_shapes(presentation.slides[source['face'] - 1])
        if len(pictures) < 2:
            raise SystemExit(f'Expected at least 2 type photos on slide {source["face"]}: {key}')
        recommended = comparison_images(presentation, source['comparison'], 'recommended')
        avoid = comparison_images(presentation, source['comparison'], 'avoid')
        fashion = slide_images(presentation, source['fashion'])
        collage_images(recommended, OUTPUT / 'female' / 'makeup' / 'recommended' / f'{key}.jpg')
        collage_images(avoid, OUTPUT / 'female' / 'makeup' / 'avoid' / f'{key}.jpg')
        collage_images(fashion, OUTPUT / 'female' / 'fashion' / f'{key}.jpg')
        card_images(recommended, OUTPUT / 'female' / 'makeup' / 'recommended' / key)
        card_images(avoid, OUTPUT / 'female' / 'makeup' / 'avoid' / key)
        card_images(fashion, OUTPUT / 'female' / 'fashion' / key)

    for group, slides in FEMALE_HAIR_RECOMMENDED_SLIDES.items():
        images = slide_images(presentation, slides)
        collage_images(images, OUTPUT / 'female' / 'hair' / 'recommended' / f'{group}.jpg')
        card_images(images, OUTPUT / 'female' / 'hair' / 'recommended' / group)
    for group, slides in FEMALE_HAIR_AVOID_SLIDES.items():
        images = slide_images(presentation, slides)
        collage_images(images, OUTPUT / 'female' / 'hair' / 'avoid' / f'{group}.jpg')
        card_images(images, OUTPUT / 'female' / 'hair' / 'avoid' / group)
    for group, slides in MALE_HAIR_RECOMMENDED_SLIDES.items():
        images = slide_images(presentation, slides)
        collage_images(images, OUTPUT / 'male' / 'hair' / f'{group}.jpg')
        card_images(images, OUTPUT / 'male' / 'hair' / group)
    for group, slides in MALE_HAIR_AVOID_SLIDES.items():
        images = slide_images(presentation, slides)
        collage_images(images, OUTPUT / 'male' / 'hair' / 'avoid-ppt' / f'{group}.jpg')
        card_images(images, OUTPUT / 'male' / 'hair' / 'avoid-ppt' / group)


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--ppt', type=Path, default=DEFAULT_PPT)
    args = parser.parse_args()
    import_assets(args.ppt)
